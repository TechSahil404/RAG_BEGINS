 
from langchain_community.document_loaders import TextLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain.vectorstores import FAISS
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM, pipeline
from langchain_community.document_loaders import UnstructuredPDFLoader
from langchain_community.document_loaders import PyMuPDFLoader


try:
    from langchain_community.document_loaders import PDFPlumberLoader
except Exception:
    PDFPlumberLoader = None

try:
    from langchain_community.document_loaders import Docx2txtLoader
except Exception:
    Docx2txtLoader = None

try:
    from langchain_community.document_loaders import PptxLoader
except Exception:
    PptxLoader = None


from langchain_community.llms.huggingface_pipeline import HuggingFacePipeline

# OCR
try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except Exception:
    pytesseract = None
    Image = None
    OCR_AVAILABLE = False

# pdf image conversion fallback
try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except Exception:
    PDF2IMAGE_AVAILABLE = False

import os
import tempfile
import streamlit as st
from typing import List
from langchain.schema import Document

# ------------------ Utility: flexible loader ------------------
def load_documents_from_path(file_path: str) -> List[Document]:
    lower = file_path.lower()
    # TXT
    if lower.endswith(".txt"):
        Loader = TextLoader(file_path, encoding="utf-8")
        return Loader.load()

    # DOCX
    if lower.endswith(".docx"):
        if Docx2txtLoader is not None:
            try:
                loader = Docx2txtLoader(file_path)
                return loader.load()
            except Exception:
                pass
        # fallback: basic read (not ideal)
        try:
            import docx
            doc = docx.Document(file_path)
            full = "\n\n".join([p.text for p in doc.paragraphs])
            return [Document(page_content=full, metadata={"source": os.path.basename(file_path)})]
        except Exception as e:
            raise RuntimeError(f"DOCX extract failed: {e}")

    # PPTX
    if lower.endswith(".pptx"):
        if PptxLoader is not None:
            try:
                loader = PptxLoader(file_path)
                return loader.load()
            except Exception:
                pass
        try:
            from pptx import Presentation
            pres = Presentation(file_path)
            slides = []
            for s in pres.slides:
                texts = []
                for shape in s.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
                slides.append("\n".join(texts))
            joined = "\n\n".join(slides)
            return [Document(page_content=joined, metadata={"source": os.path.basename(file_path)})]
        except Exception as e:
            raise RuntimeError(f"PPTX extract failed: {e}")

    # Images -> OCR
    if lower.endswith((".png", ".jpg", ".jpeg", ".tiff", ".bmp")):
        if not OCR_AVAILABLE:
            raise RuntimeError("OCR not available (pytesseract/PIL missing).")
        try:
            text = pytesseract.image_to_string(Image.open(file_path))
            return [Document(page_content=text, metadata={"source": os.path.basename(file_path)})]
        except Exception as e:
            raise RuntimeError(f"OCR failed: {e}")

    # PDF: try multiple approaches
    if lower.endswith(".pdf"):
        # 1) Unstructured
        try:
            loader = UnstructuredPDFLoader(file_path, mode="elements")
            docs = loader.load()
            if docs and any(d.page_content.strip() for d in docs):
                return docs
        except Exception:
            pass

        # 2) PyMuPDF
        try:
            loader = PyMuPDFLoader(file_path)
            docs = loader.load()
            if docs and any(d.page_content.strip() for d in docs):
                return docs
        except Exception:
            pass

        # 3) pdfplumber loader if available
        if PDFPlumberLoader is not None:
            try:
                loader = PDFPlumberLoader(file_path)
                docs = loader.load()
                if docs and any(d.page_content.strip() for d in docs):
                    return docs
            except Exception:
                pass

        # 4) OCR fallback (pdf2image + pytesseract)
        if OCR_AVAILABLE and PDF2IMAGE_AVAILABLE:
            try:
                pages = convert_from_path(file_path)
                texts = []
                for p in pages:
                    texts.append(pytesseract.image_to_string(p))
                joined = "\n\n".join(texts)
                return [Document(page_content=joined, metadata={"source": os.path.basename(file_path)})]
            except Exception:
                pass

        raise RuntimeError("Failed to extract text from PDF with available loaders/OCR on this system.")

    raise RuntimeError("Unsupported file type: " + file_path)

# ------------------ Chunking ------------------
def create_chunks(docs: List[Document], chunk_size=1000, chunk_overlap=200) -> List[Document]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", " ", ""]
    )
    return splitter.split_documents(docs)

# ------------------ Embeddings & Vectorstore ------------------
def init_embeddings(model_name="sentence-transformers/all-mpnet-base-v2"):
    return HuggingFaceEmbeddings(model_name=model_name)

def build_vectorstore(chunks: List[Document], embeddings_model, persist_path: str = None):
    vs = FAISS.from_documents(documents=chunks, embedding=embeddings_model)
    if persist_path:
        vs.save_local(persist_path)
    return vs

# ------------------ LLM Setup ------------------
def init_llm(model_name="google/flan-t5-small", max_new_tokens=200):
    tokenizer = AutoTokenizer.from_pretrained(model_name)
    model = AutoModelForSeq2SeqLM.from_pretrained(model_name)
    pipe = pipeline(
        "text2text-generation",
        model=model,
        tokenizer=tokenizer,
        max_new_tokens=max_new_tokens
    )
    llm = HuggingFacePipeline(pipeline=pipe)
    return llm

def parse_llm_output(raw):
    # handle various pipeline outputs
    try:
        if isinstance(raw, list) and len(raw) > 0:
            first = raw[0]
            if isinstance(first, dict):
                for k in ("generated_text", "text", "generated_texts"):
                    if k in first:
                        return first[k]
                # fallback
                return list(first.values())[0]
            else:
                return str(first)
        return str(raw)
    except Exception:
        return str(raw)

# ------------------ RAG High-level ------------------
def answer_query_with_rag(query: str, vector_store: FAISS, llm, k=3):
    retrieved = vector_store.similarity_search(query, k=k)
    context = "\n\n".join([d.page_content for d in retrieved])
    prompt = (
        "Answer the question based on the context below. If the answer is not present in the context, say 'I don't know'.\n\n"
        f"Context:\n{context}\n\nQuestion: {query}\n\nAnswer:"
    )
    raw = llm(prompt)
    ans = parse_llm_output(raw)
    return ans, retrieved

# ------------------ STREAMLIT UI ------------------
st.set_page_config(page_title="Simple RAG Multi-format", layout="wide")
st.title("📚 Simple RAG Demo — Upload docs & ask questions")

st.markdown("""
**Supported:** PDF (unstructured + fallback), TXT, DOCX, PPTX, Images (OCR).  
**Notes:** OCR requires Tesseract + pytesseract. PDF image fallback requires poppler + pdf2image.
""")

# session state
if "vector_store" not in st.session_state:
    st.session_state.vector_store = None
if "embeddings_model" not in st.session_state:
    st.session_state.embeddings_model = None
if "llm" not in st.session_state:
    st.session_state.llm = None
if "chunks" not in st.session_state:
    st.session_state.chunks = None
if "uploaded_files" not in st.session_state:
    st.session_state.uploaded_files = []

# sidebar
with st.sidebar:
    st.header("Settings")
    preload = st.checkbox("Init models on start (may take time)", value=False)
    chunk_size = st.number_input("Chunk size", min_value=200, max_value=4000, value=1000, step=100)
    chunk_overlap = st.number_input("Chunk overlap", min_value=50, max_value=2000, value=200, step=50)
    k = st.number_input("Retriever k", min_value=1, max_value=10, value=3, step=1)
    emb_model_name = st.selectbox("Embeddings model", ["sentence-transformers/all-mpnet-base-v2"], index=0)
    llm_model_name = st.selectbox("LLM (CPU)", ["google/flan-t5-small"], index=0)

    if preload and st.button("Init models now"):
        with st.spinner("Initializing embeddings and LLM..."):
            st.session_state.embeddings_model = init_embeddings(model_name=emb_model_name)
            st.session_state.llm = init_llm(model_name=llm_model_name)
        st.success("Models initialized.")

# File uploader
uploaded = st.file_uploader("Upload files (pdf/txt/docx/pptx/png/jpg)", accept_multiple_files=True)

if uploaded:
    st.session_state.uploaded_files = uploaded
    all_docs = []
    errors = []
    for up in uploaded:
        suffix = os.path.splitext(up.name)[1]
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        tmp.write(up.getbuffer())
        tmp.flush()
        tmp.close()
        st.info(f"Processing {up.name} ...")
        try:
            docs = load_documents_from_path(tmp.name)
            for d in docs:
                if hasattr(d, "metadata"):
                    d.metadata["source"] = up.name
            all_docs.extend(docs)
            st.success(f"Loaded {len(docs)} documents from {up.name}")
        except Exception as e:
            errors.append((up.name, str(e)))
            st.error(f"Failed to load {up.name}: {e}")

    if len(all_docs) == 0:
        st.warning("No textual content extracted from uploaded files. Check errors above.")
    else:
        st.info("Chunking documents ...")
        chunks = create_chunks(all_docs, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
        st.session_state.chunks = chunks
        st.success(f"Created {len(chunks)} chunks.")

        if st.session_state.embeddings_model is None:
            with st.spinner("Initializing embeddings ..."):
                st.session_state.embeddings_model = init_embeddings(model_name=emb_model_name)
            st.success("Embeddings ready.")

        with st.spinner("Building FAISS index ..."):
            st.session_state.vector_store = build_vectorstore(chunks, st.session_state.embeddings_model)
        st.success("Vector store ready (FAISS).")
        try:
            st.write(f"FAISS index size: {st.session_state.vector_store.index.ntotal}")
        except Exception:
            pass

# Query UI
st.markdown("---")
query = st.text_input("❓ Enter your question (will use uploaded docs):")
if st.button("Submit Question"):
    if not query.strip():
        st.warning("Please enter a question.")
    elif st.session_state.vector_store is None:
        st.warning("No index found. Upload docs and build index first.")
    else:
        if st.session_state.llm is None:
            with st.spinner("Loading LLM (CPU) ..."):
                st.session_state.llm = init_llm(model_name=llm_model_name)
            st.success("LLM loaded.")

        with st.spinner("Retrieving and generating ..."):
            answer, retrieved = answer_query_with_rag(query, st.session_state.vector_store, st.session_state.llm, k=k)
        st.markdown("### ✅ Answer")
        st.write(answer)

        st.markdown("### 🔎 Retrieved Chunks")
        for i, d in enumerate(retrieved):
            src = d.metadata.get("source", "unknown")
            st.write(f"**Chunk {i+1}** — source: {src}")
            snippet = d.page_content[:800] + ("..." if len(d.page_content) > 800 else "")
            st.code(snippet)

st.markdown("---")
st.caption("Tip: FLAN-T5-small is CPU-friendly but limited in quality. For better answers use a stronger model or hosted LLM.")

# End of app
